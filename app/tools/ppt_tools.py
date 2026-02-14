"""
PowerPoint analysis tools
"""
import json
from typing import Dict, Any, List
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTAnalyzer:
    """Tools for analyzing PowerPoint presentations"""

    def extract_text_from_ppt(self, ppt_path: str) -> Dict[str, Any]:
        """
        Extract all text content from PowerPoint file

        Args:
            ppt_path: Path to PPT/PPTX file

        Returns:
            Dictionary with extracted content
        """
        try:
            prs = Presentation(ppt_path)

            extraction = {
                "slide_count": len(prs.slides),
                "slides": [],
                "all_text": "",
                "title_slides": [],
                "bullet_points": [],
                "notes": []
            }

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "slide_number": slide_num,
                    "title": "",
                    "content": [],
                    "notes": "",
                    "has_images": False,
                    "has_tables": False,
                    "has_charts": False
                }

                # Extract text from shapes
                for shape in slide.shapes:
                    # Check shape type
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        slide_data["has_images"] = True
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        slide_data["has_tables"] = True
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        slide_data["has_charts"] = True

                    # Extract text
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()

                        # Check if it's a title
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            placeholder = shape.placeholder_format
                            if placeholder.type == 1:  # Title placeholder
                                slide_data["title"] = text
                                extraction["title_slides"].append({
                                    "slide_number": slide_num,
                                    "title": text
                                })
                                continue

                        # Regular content
                        if text:
                            slide_data["content"].append(text)
                            extraction["all_text"] += text + "\n"

                            # Detect bullet points
                            if text.startswith("•") or text.startswith("-") or text.startswith("*"):
                                extraction["bullet_points"].append({
                                    "slide": slide_num,
                                    "text": text
                                })

                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text:
                        slide_data["notes"] = notes_text
                        extraction["notes"].append({
                            "slide": slide_num,
                            "notes": notes_text
                        })

                extraction["slides"].append(slide_data)

            return extraction

        except Exception as e:
            return {"error": f"Error extracting PPT content: {str(e)}"}

    def analyze_presentation_structure(self, ppt_path: str) -> Dict[str, Any]:
        """
        Analyze presentation structure and content organization

        Args:
            ppt_path: Path to PPT/PPTX file

        Returns:
            Dictionary with structural analysis
        """
        extraction = self.extract_text_from_ppt(ppt_path)

        if "error" in extraction:
            return extraction

        analysis = {
            "slide_count": extraction["slide_count"],
            "has_title_slide": len(extraction["title_slides"]) > 0,
            "sections_identified": self._identify_sections(extraction),
            "content_distribution": self._analyze_content_distribution(extraction),
            "visual_elements": self._count_visual_elements(extraction),
            "key_slides": self._identify_key_slides(extraction),
            "extracted_text": extraction["all_text"][:5000]  # First 5000 chars for LLM
        }

        return analysis

    def _identify_sections(self, extraction: Dict[str, Any]) -> List[Dict[str, str]]:
        """Identify main sections from slide titles"""
        sections = []
        for title_slide in extraction["title_slides"]:
            title = title_slide["title"].lower()

            # Common section keywords
            section_type = "other"
            if any(word in title for word in ["problem", "challenge", "issue"]):
                section_type = "problem"
            elif any(word in title for word in ["solution", "our approach", "how it works"]):
                section_type = "solution"
            elif any(word in title for word in ["market", "opportunity", "business"]):
                section_type = "market"
            elif any(word in title for word in ["team", "about us"]):
                section_type = "team"
            elif any(word in title for word in ["demo", "product"]):
                section_type = "demo"
            elif any(word in title for word in ["competitive", "competition", "vs"]):
                section_type = "competition"

            sections.append({
                "slide": title_slide["slide_number"],
                "title": title_slide["title"],
                "type": section_type
            })

        return sections

    def _analyze_content_distribution(self, extraction: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze how content is distributed across slides"""
        distribution = {
            "text_heavy_slides": 0,
            "visual_heavy_slides": 0,
            "balanced_slides": 0,
            "average_text_per_slide": 0
        }

        total_text_length = 0

        for slide in extraction["slides"]:
            text_length = sum(len(text) for text in slide["content"])
            total_text_length += text_length

            has_visuals = slide["has_images"] or slide["has_charts"] or slide["has_tables"]

            if text_length > 500 and not has_visuals:
                distribution["text_heavy_slides"] += 1
            elif has_visuals and text_length < 200:
                distribution["visual_heavy_slides"] += 1
            else:
                distribution["balanced_slides"] += 1

        if extraction["slide_count"] > 0:
            distribution["average_text_per_slide"] = total_text_length / extraction["slide_count"]

        return distribution

    def _count_visual_elements(self, extraction: Dict[str, Any]) -> Dict[str, int]:
        """Count visual elements across all slides"""
        visuals = {
            "images": 0,
            "tables": 0,
            "charts": 0
        }

        for slide in extraction["slides"]:
            if slide["has_images"]:
                visuals["images"] += 1
            if slide["has_tables"]:
                visuals["tables"] += 1
            if slide["has_charts"]:
                visuals["charts"] += 1

        return visuals

    def _identify_key_slides(self, extraction: Dict[str, Any]) -> List[int]:
        """Identify potentially important slides"""
        key_slides = []

        for slide in extraction["slides"]:
            slide_num = slide["slide_number"]
            title = slide["title"].lower()
            content = " ".join(slide["content"]).lower()

            # Check for important keywords
            important_keywords = [
                "problem", "solution", "market", "revenue", "business model",
                "competitive", "advantage", "unique", "innovation",
                "results", "metrics", "growth", "traction"
            ]

            if any(keyword in title or keyword in content for keyword in important_keywords):
                if slide_num not in key_slides:
                    key_slides.append(slide_num)

        return sorted(key_slides)


# Convenience function for CrewAI tools
def analyze_ppt(ppt_path: str) -> str:
    """
    Analyze PowerPoint presentation (for use as CrewAI tool)

    Args:
        ppt_path: Path to PPT/PPTX file

    Returns:
        str: JSON string with analysis
    """
    analyzer = PPTAnalyzer()

    # Get detailed extraction
    extraction = analyzer.extract_text_from_ppt(ppt_path)

    # Get structural analysis
    structure = analyzer.analyze_presentation_structure(ppt_path)

    result = {
        "content_extraction": extraction,
        "structural_analysis": structure
    }

    return json.dumps(result, indent=2)
